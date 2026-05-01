"""Build the HypothesisLoop presentation deck from the outline.

Outputs: docs/HypothesisLoop_Presentation.pptx

Mission-control aesthetic: slate background (#0F172A), panel (#1E293B),
cyan accent (#06B6D4), JetBrains Mono / Consolas typography. Mirrors the
Streamlit UI's look so the deck and the live demo feel like one product.

Re-run this after editing the outline to regenerate.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---- Mission-control palette ----
BG = RGBColor(0x0F, 0x17, 0x2A)
PANEL = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0x33, 0x41, 0x55)
FG = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xF8, 0x71, 0x71)
AMBER = RGBColor(0xFB, 0xBF, 0x24)

MONO = "Consolas"  # ships with Windows; JetBrains Mono if installed will look better
SANS = "Segoe UI"

OUT_PATH = Path(__file__).resolve().parents[1] / "docs" / "HypothesisLoop_Presentation.pptx"


def _set_slide_background(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=FG, font=MONO, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _add_panel(slide, left, top, width, height, *, fill=PANEL, line=BORDER, line_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    # Kill the default outline + remove text
    shape.text_frame.text = ""
    return shape


def _add_accent_bar(slide, left, top, width, height, color=ACCENT):
    """Thin vertical or horizontal accent bar — used for left-border emphasis."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_topbar(slide, slide_num, total_slides, section_label):
    """Mission-control top bar with brand left + slide-num + section right."""
    top_h = Inches(0.6)
    _add_panel(slide, 0, 0, Inches(13.333), top_h, fill=PANEL, line=BORDER)
    _add_textbox(
        slide, Inches(0.4), Inches(0.1), Inches(5), Inches(0.4),
        "◐ HYPOTHESISLOOP",
        size=14, bold=True, color=ACCENT, font=MONO, anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide, Inches(7.5), Inches(0.1), Inches(3), Inches(0.4),
        f"SECTION  {section_label.upper()}",
        size=10, color=MUTED, font=MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT,
    )
    _add_textbox(
        slide, Inches(10.6), Inches(0.1), Inches(2.4), Inches(0.4),
        f"SLIDE  {slide_num:02d} / {total_slides:02d}",
        size=10, color=MUTED, font=MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT,
    )


def _slide_title(slide, text):
    """Big section heading below the topbar."""
    _add_accent_bar(slide, Inches(0.4), Inches(1.0), Inches(0.05), Inches(0.6), color=ACCENT)
    _add_textbox(
        slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6),
        text,
        size=28, bold=True, color=FG, font=SANS, anchor=MSO_ANCHOR.MIDDLE,
    )


def _bullet_block(slide, left, top, width, height, items, *, size=16, gap_pt=10):
    """A list of bullets; each item is a string. Starts with a cyan ►."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else gap_pt)
        p.space_after = Pt(0)

        run_arrow = p.add_run()
        run_arrow.text = "▸ "
        run_arrow.font.name = MONO
        run_arrow.font.size = Pt(size)
        run_arrow.font.color.rgb = ACCENT
        run_arrow.font.bold = True

        run_text = p.add_run()
        run_text.text = item
        run_text.font.name = SANS
        run_text.font.size = Pt(size)
        run_text.font.color.rgb = FG


# ---- Slide builders ----

def slide_title(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "TITLE")

    # Big brand-style title centered
    _add_textbox(
        s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.0),
        "HypothesisLoop",
        size=72, bold=True, color=ACCENT, font=SANS, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.6),
        "An LLM-Powered Iterative Hypothesis-Testing Agent",
        size=22, color=FG, font=SANS, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s, Inches(0.6), Inches(3.9), Inches(12), Inches(0.4),
        "for Tabular Data Analysis",
        size=22, color=FG, font=SANS, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Author line in mono
    _add_textbox(
        s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
        "HOLT  YOUNG    +    SAM  PENN",
        size=16, color=ACCENT, font=MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.4),
        "QAC 387 — Wesleyan — Spring 2026",
        size=13, color=MUTED, font=MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )


def slide_problem(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Problem")
    _slide_title(s, "The Problem")

    _bullet_block(s, Inches(0.7), Inches(2.0), Inches(12), Inches(4.5), [
        "ChatGPT Code Interpreter is one-shot — you ask, it answers, you figure out the next question yourself.",
        "Real data analysis is iterative: hypothesize → test → learn → repeat. Most people lack the stats/coding background to do that loop unaided.",
        "Existing LLM tools answer single queries; they don't follow up on what they find.",
        "We wanted an agent that wraps the scientific method around a CSV and runs it for the user.",
    ], size=17, gap_pt=14)


def slide_what_we_built(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "What we built")
    _slide_title(s, "What We Built")

    _bullet_block(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.0), [
        "Python LLM agent that runs the scientific method loop on a CSV:  Hypothesize → Experiment → Evaluate → Learn → Repeat.",
        "Default 5 iterations, fully autonomous, every step traced in Langfuse.",
    ], size=17, gap_pt=12)

    # Two-mode panel
    _add_panel(s, Inches(0.7), Inches(4.2), Inches(5.8), Inches(2.5), fill=PANEL, line=BORDER)
    _add_accent_bar(s, Inches(0.7), Inches(4.2), Inches(0.05), Inches(2.5), color=ACCENT)
    _add_textbox(s, Inches(0.95), Inches(4.35), Inches(5.4), Inches(0.4),
                 "MODE  ·  EXPLORE", size=11, color=ACCENT, font=MONO, bold=True)
    _add_textbox(s, Inches(0.95), Inches(4.75), Inches(5.4), Inches(0.5),
                 "Free-form research question →", size=15, color=FG, font=SANS)
    _add_textbox(s, Inches(0.95), Inches(5.10), Inches(5.4), Inches(0.5),
                 "narrative report with charts.", size=15, color=FG, font=SANS)
    _add_textbox(s, Inches(0.95), Inches(5.7), Inches(5.4), Inches(0.8),
                 "Use case: \"tell me what's interesting in this data.\"",
                 size=12, color=MUTED, font=MONO)

    _add_panel(s, Inches(6.85), Inches(4.2), Inches(5.8), Inches(2.5), fill=PANEL, line=BORDER)
    _add_accent_bar(s, Inches(6.85), Inches(4.2), Inches(0.05), Inches(2.5), color=EMERALD)
    _add_textbox(s, Inches(7.10), Inches(4.35), Inches(5.4), Inches(0.4),
                 "MODE  ·  PREDICT", size=11, color=EMERALD, font=MONO, bold=True)
    _add_textbox(s, Inches(7.10), Inches(4.75), Inches(5.4), Inches(0.5),
                 "Pick target column → engineered", size=15, color=FG, font=SANS)
    _add_textbox(s, Inches(7.10), Inches(5.10), Inches(5.4), Inches(0.5),
                 "features + AutoGluon ensemble.", size=15, color=FG, font=SANS)
    _add_textbox(s, Inches(7.10), Inches(5.7), Inches(5.4), Inches(0.8),
                 "Use case: \"build me the best predictor of column X.\"",
                 size=12, color=MUTED, font=MONO)


def slide_tech_stack(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Tech stack")
    _slide_title(s, "Tech Stack")

    rows = [
        ("LLMs",         "Moonshot Kimi K2.6 (default, ~10× cheaper than Opus)  ·  OpenAI GPT-4o-mini fallback"),
        ("Embeddings",   "OpenAI text-embedding-3-small (novelty detection, 1,536-dim, cosine gate)"),
        ("Vector store", "FAISS (local, in-memory, codebook + statistical-test guide)"),
        ("Tracing",      "Langfuse — every LLM call, every iteration, every retry"),
        ("ML",           "pandas · scipy · sklearn · AutoGluon (ensemble for Predict mode)"),
        ("UI",           "Streamlit + custom CSS  ·  CLI with --auto / --resume / --report-only"),
        ("Sandbox",      "subprocess + AST denylist + RLIMIT_AS + 30s timeout (no Docker)"),
        ("Cost",         "~$0.05–$0.15 per run  ·  ~$1.30 for the 28-run validation suite"),
    ]
    top = Inches(2.0)
    row_h = Inches(0.55)
    for i, (k, v) in enumerate(rows):
        y = top + row_h * i
        _add_textbox(s, Inches(0.8), y, Inches(2.2), row_h,
                     k.upper(), size=12, color=ACCENT, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(s, Inches(3.1), y, Inches(9.8), row_h,
                     v, size=14, color=FG, font=SANS, anchor=MSO_ANCHOR.MIDDLE)


def slide_loop_works(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "How the loop works")
    _slide_title(s, "How the Loop Actually Works")

    # 5 step cards
    steps = [
        ("HYPOTHESIZE", "LLM proposes a falsifiable claim.", "In Predict mode, commits to a metric Δ."),
        ("EXPERIMENT",  "LLM writes Python.",                "Sandbox runs it. Retries on error (≤4 attempts)."),
        ("EVALUATE",    "LLM interprets output.",            "In Predict, deterministic accept/reject overrides."),
        ("LEARN",       "DAG-tracked trace updated.",        "Novelty gate, soft-decay, bias scanner."),
        ("REPEAT",      "Next iteration sees prior",         "hypotheses, decisions, AND code failures."),
    ]
    top = Inches(2.0)
    h = Inches(4.5)
    n_steps = len(steps)
    total_w = Inches(12.5)
    gap = Inches(0.1)
    card_w = (total_w - gap * (n_steps - 1)) / n_steps

    for i, (label, line1, line2) in enumerate(steps):
        x = Inches(0.4) + (card_w + gap) * i
        _add_panel(s, x, top, card_w, h, fill=PANEL, line=BORDER)
        _add_accent_bar(s, x, top, card_w, Inches(0.06), color=ACCENT)
        _add_textbox(s, x + Inches(0.15), top + Inches(0.25), card_w - Inches(0.3), Inches(0.5),
                     label, size=12, color=ACCENT, font=MONO, bold=True)
        _add_textbox(s, x + Inches(0.15), top + Inches(0.85), card_w - Inches(0.3), Inches(2.5),
                     line1, size=13, color=FG, font=SANS)
        _add_textbox(s, x + Inches(0.15), top + Inches(2.35), card_w - Inches(0.3), Inches(2.0),
                     line2, size=11, color=MUTED, font=SANS)


def slide_pivots(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Iterations & pivots")
    _slide_title(s, "Key Pivots")

    _add_textbox(s, Inches(0.7), Inches(1.85), Inches(12.0), Inches(0.45),
                 "We rebuilt mid-project. Here's what changed and why.",
                 size=15, color=MUTED, font=SANS)

    pivots = [
        ("Build 4 → HypothesisLoop",
         "Abandoned the cookie-cutter REPL agent and rebuilt around the original proposal's autonomous loop."),
        ("Opus 4.6 → Kimi K2.6",
         "Same tabular-reasoning quality at ~1/10th the cost. 128K context handles full traces without aggressive pruning."),
        ("sklearn baseline → AutoGluon",
         "Performance: AutoGluon ranks #1 on OpenML AutoML benchmark (63% top-1) — leaves significant accuracy on the table to use anything else."),
        ("Hard novelty ban → layered gate",
         "Prompt-injection + embedding cosine 0.85 + soft-decay (loosens to 0.92 after 3 consecutive rejections)."),
        ("Per-iteration retry only → cross-iteration error context",
         "Next iteration's hypothesize prompt sees prior iterations' code crashes, so the agent stops repeating the same dtype mistakes."),
    ]
    top = Inches(2.45)
    for i, (head, body) in enumerate(pivots):
        y = top + Inches(0.8) * i
        _add_textbox(s, Inches(0.8), y, Inches(11.8), Inches(0.4),
                     head.upper(), size=11, color=ACCENT, font=MONO, bold=True)
        _add_textbox(s, Inches(0.8), y + Inches(0.3), Inches(11.8), Inches(0.5),
                     body, size=14, color=FG, font=SANS)


def slide_bias_safety(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Bias scanner & safety")
    _slide_title(s, "Bias Scanner & Safety")

    _bullet_block(s, Inches(0.8), Inches(2.0), Inches(12), Inches(3.5), [
        "Pattern: sensitive variable (race · sex · native-country · marital-status) + causal verb (causes · leads to · due to · drives) → flag.",
        "Flagged paragraphs get a ⚠ correlational-only disclaimer prepended in the report; live UI shows a red chip on the iteration card.",
        "Question-level scan (added late): even when the LLM self-disciplines on output, a causal user prompt still triggers a top-of-timeline banner.",
        "Sandbox AST denylist: blocks os, subprocess, eval, file writes outside session dir, target-column references in feature engineering.",
        "ASCII identifier lint: rejects non-ASCII identifiers (e.g. CJK from hallucinations) before sandbox execution.",
    ], size=15, gap_pt=10)

    # Disclaimer pill at bottom
    _add_panel(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9),
               fill=RGBColor(0x3F, 0x1F, 0x1F), line=RED)
    _add_textbox(s, Inches(1.0), Inches(6.05), Inches(11.5), Inches(0.8),
                 "⚠  Causal claim about a sensitive variable — interpret as correlational only.",
                 size=14, bold=True, color=RGBColor(0xFE, 0xCA, 0xCA), font=MONO, anchor=MSO_ANCHOR.MIDDLE)


def slide_validation(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Validation")
    _slide_title(s, "Validation — 28 Runs Across 12 Categories")

    rows = [
        ("Total runs",                       "28 / 28 completed",                 "✓"),
        ("Wall time",                        "120.6 minutes",                     "·"),
        ("Cost",                             "~$1.30 (vs. $25 cap)",              "✓"),
        ("Loop completion (no crash)",       "89% (25 / 28)",                     "⚠"),
        ("Codegen success (raw exit 0)",     "45%",                               "⚠"),
        ("Codegen success (with retries)",   "87%",                               "✓"),
        ("Novelty gate firing",              "28 rejections across 149 iters",    "✓"),
        ("Bias scanner — node level",        "0 / 0 fires (LLM self-disciplines)", "✓"),
        ("Bias scanner — question level",    "Catches every causal-framed prompt", "✓"),
        ("Final report renders",             "27 / 28 (96%)",                     "✓"),
        ("Predict-mode AutoGluon lift",      "Baseline 0.905 → Final 0.932 ROC AUC", "✓"),
    ]
    top = Inches(2.0)
    row_h = Inches(0.42)
    for i, (k, v, status) in enumerate(rows):
        y = top + row_h * i
        if i % 2 == 0:
            _add_panel(s, Inches(0.6), y, Inches(12.1), row_h, fill=PANEL, line=PANEL)
        c = EMERALD if status == "✓" else (AMBER if status == "⚠" else MUTED)
        _add_textbox(s, Inches(0.8), y, Inches(0.35), row_h,
                     status, size=14, color=c, font=MONO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(s, Inches(1.2), y, Inches(6.5), row_h,
                     k, size=13, color=FG, font=SANS, anchor=MSO_ANCHOR.MIDDLE)
        _add_textbox(s, Inches(7.7), y, Inches(5.0), row_h,
                     v, size=13, color=ACCENT, font=MONO, anchor=MSO_ANCHOR.MIDDLE)


def slide_limitations(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Limitations")
    _slide_title(s, "Honest Limitations")

    _bullet_block(s, Inches(0.8), Inches(2.0), Inches(12), Inches(4.5), [
        "Windows can't enforce RAM cap (POSIX-only setrlimit). Documented; timeout is the primary defense.",
        "Kimi nondeterminism: ~50% of attempts crash on a trailing plot block even when the analysis itself worked. Retry catches most; metrics still emit.",
        "Predict mode's +0.001 ROC-AUC threshold rejects most LLM-proposed features. Conservatism is intentional but limits visible feature acceptance in short runs.",
        "Stateful FE transforms (target encoding, mean-centering) silently degrade quality on the test split. Mitigated by prompt rules; real fix is sklearn fit/transform pipelines.",
        "Bias scanner has zero hit rate at the node level — the LLM self-disciplines on causal language. Question-level scan added as backstop.",
    ], size=15, gap_pt=11)


def slide_next_steps(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Next steps")
    _slide_title(s, "Next Steps")

    _bullet_block(s, Inches(0.8), Inches(2.0), Inches(12), Inches(4.5), [
        "sklearn fit/transform pipeline for stateful FE (the proper fix for the test-split quality drift).",
        "Multi-target prediction; multi-objective optimization across competing metrics.",
        "Persistent embedding cache (currently in-process LRU only).",
        "Cross-dataset benchmarks: House Prices (regression), WHO Life Expectancy.",
        "Schedule-driven scheduler — replace LinearScheduler with RD-Agent-style probabilistic exploration.",
        "Hyphenated target columns (P2.1 crashed on `hours-per-week`); plain wiring fix.",
        "Live mid-iteration sub-step status (currently single-string label until iteration completes).",
    ], size=15, gap_pt=10)


def slide_demo(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Live demo")
    _slide_title(s, "Live Demo")

    # Big mono "TERMINAL" treatment
    _add_panel(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.5), fill=PANEL, line=BORDER)
    _add_accent_bar(s, Inches(1.0), Inches(2.2), Inches(0.06), Inches(4.5), color=ACCENT)
    lines = [
        "$ python -m streamlit run hypothesisloop/ui/streamlit_app.py",
        "",
        "▸ Upload data/adult.csv",
        "▸ Switch to PREDICT mode",
        "▸ Target = income  ·  task = classification (auto)",
        "▸ Auto-run · max-iters 5 · AutoML budget 120s",
        "▸ Hit ▶ START RUN",
        "▸ Watch iterations stream into the timeline",
        "▸ Use ↻ CONTINUE +5 ITERS to extend the run mid-flight",
        "▸ Open the generated report.md and AutoGluon leaderboard.csv",
    ]
    _add_textbox(s, Inches(1.4), Inches(2.45), Inches(10.9), Inches(4.2),
                 "\n".join(lines), size=15, color=FG, font=MONO)

    _add_textbox(s, Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.4),
                 "Backup: pre-recorded run in reports/hl-20260501-045426-1a98 — same flow, same artifacts.",
                 size=11, color=MUTED, font=MONO, anchor=MSO_ANCHOR.MIDDLE)


def slide_takeaways(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(s, BG)
    _add_topbar(s, n, total, "Takeaways")
    _slide_title(s, "Takeaways")

    _bullet_block(s, Inches(0.8), Inches(2.0), Inches(12), Inches(4.0), [
        "We pivoted away from the assignment scaffolding and built what we actually proposed.",
        "Auditable: every hypothesis, every line of code, every decision lives in Langfuse and a Markdown report.",
        "Two modes from one loop: insights AND a deployable model from the same engine.",
        "~$0.05 / run.  100% local sandbox.  No vendor lock-in beyond the LLM provider.",
        "Predict mode shipped a real ROC-AUC lift: 0.905 baseline → 0.932 AutoGluon ensemble on UCI Adult.",
    ], size=16, gap_pt=14)

    _add_textbox(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.5),
                 "Q & A", size=24, bold=True, color=ACCENT, font=SANS, align=PP_ALIGN.CENTER)


# ---- Build ----

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title,
        slide_problem,
        slide_what_we_built,
        slide_tech_stack,
        slide_loop_works,
        slide_pivots,
        slide_bias_safety,
        slide_validation,
        slide_limitations,
        slide_next_steps,
        slide_demo,
        slide_takeaways,
    ]
    total = len(builders)
    for i, b in enumerate(builders, 1):
        b(prs, i, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}  ({OUT_PATH.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    build()
